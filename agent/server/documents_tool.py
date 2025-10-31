"""
mcp_document_processing.py
FastMCP server exposing a Document‑processing tool
that works without the camel package.

File Reading Strategy:
All files for processing should be placed in the unified workspace directory 
(task_cache_dir/workspace/) as managed by agent_core.py. This tool receives
the workspace directory path and filename as separate parameters, then combines
them for file operations to maintain consistency with the code execution environment.
"""

# --------------------------------------------------------------------------- #
#  Imports
# --------------------------------------------------------------------------- #
import asyncio, os, io, json, subprocess
from typing import Tuple, Optional, List, Literal
from pathlib import Path

from loguru import logger
from retry import retry

from mcp.server.fastmcp import FastMCP
import anyio

# --- your own helper toolkits ------------------------------------------------ #
#   (provide these scripts in the Python path)
from image_tool import ask_question_about_image
from excel_tool import ExcelToolkit
try:
    from video_tool import ask_question_about_video
except Exception as import_error:
    logger.warning(f"Video tool unavailable, disabling video support in documents_tool: {import_error}")

    def ask_question_about_video(*args, **kwargs):
        raise RuntimeError(
            "Video processing is unavailable because the video tool backend "
            "could not be loaded (missing dependency such as OpenCV/libGL)."
        )
# --- third‑party libs already used ------------------------------------------ #
import assemblyai as aai
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from PIL import Image
from docx2markdown._docx_to_markdown import docx_to_markdown
from chunkr_ai import Chunkr
import xmltodict
import nest_asyncio
nest_asyncio.apply()

from dotenv import load_dotenv
load_dotenv(".env")


# --------------------------------------------------------------------------- #
#  Unified path management functions
# --------------------------------------------------------------------------- #
def _get_workspace_dir(task_cache_dir: Optional[str] = None) -> str:
    """
    Get or create the unified workspace directory for a task.
    This follows the same pattern as code_tool.py to ensure consistency.
    
    📁 DIRECTORY STRUCTURE:
        task_cache_dir/workspace/  (flat structure - all files here)
    
    Args:
        task_cache_dir: Task-specific cache directory path
        
    Returns:
        str: Absolute path to workspace directory
    """
    if not task_cache_dir:
        task_cache_dir = os.getenv("AGENT_CACHE_DIR", "./agent_cache")
    
    workspace_dir = Path(task_cache_dir) / "workspace"
    workspace_dir.mkdir(parents=True, exist_ok=True)
    return str(workspace_dir)


def _build_file_path(workspace_dir: str, filename: str) -> str:
    """
    Builds a secure and correct file path within the workspace.

    This function ensures that the requested file path is safely contained
    within the designated workspace directory, preventing directory traversal
    attacks (e.g., '../../etc/passwd'). It correctly handles filenames that
    include subdirectories, such as 'upload_files/document.pdf'.

    Args:
        workspace_dir: The absolute path to the task's workspace.
        filename: The relative path of the file from the workspace root.

    Returns:
        A securely resolved, absolute path to the requested file.
    
    Raises:
        ValueError: If the resolved path attempts to escape the workspace directory.
    """
    # Create absolute paths for security checks
    workspace_path = Path(workspace_dir).resolve()
    # Combine the workspace path with the relative filename
    full_path = (workspace_path / filename).resolve()

    # Security check: Ensure the resolved path is still inside the workspace
    if not str(full_path).startswith(str(workspace_path)):
        raise ValueError(
            f"Security Error: Attempted file access outside of workspace. "
            f"Original: '{filename}', Resolved: '{full_path}'"
        )
    
    return str(full_path)


# --------------------------------------------------------------------------- #
#  Toolkit implementation (no camel.BaseToolkit!)
# --------------------------------------------------------------------------- #
class DocumentProcessingToolkit:
    """
    This tool exposes a **general‑purpose document‑processing endpoint** that
    converts almost any common file you point it to into **clean, readable
    text or Markdown**.  It is useful whenever an agent needs to "look inside"
    an arbitrary file before reasoning over its contents.

    **WORKSPACE INTEGRATION**: 
    All files for processing must be located in the unified workspace directory
    (task_cache_dir/workspace/) as established by the code execution environment.
    This ensures consistent file access across all tools in the agent system.

    Conceptually, you can think of it as an *all‑in‑one* "open the file and
    give me the text" utility:

    • **Images (.jpg / .jpeg / .png)** – runs a vision model and returns a
    detailed caption.  
    • **Audio (.mp3 / .wav / .m4a)** – performs automatic transcription.  
    • **PowerPoint (.pptx)** – pulls every textbox, captions each embedded
    image, and preserves the slide order.  
    • **Spreadsheets (.xls / .xlsx / .csv)** – dumps cell values in a
    readable, row‑wise layout.  
    • **ZIP archives** – unpacks the archive and lists the extracted files.  
    • **Plain text‑like formats (.py / .txt)** – simply reads the file.  
    • **JSON, JSONL, JSON‑LD** – returns the parsed JSON structure.  
    • **XML** – converts to a Python dict (falls back to raw XML on error).  
    • **Word (.docx)** – converts the entire document to Markdown.  
    • **Video (.mov)** – generates a descriptive narration of the clip.  
    • **PDF or any other format** – attempts Chunkr AI extraction first, and
    then a plain‑text PDF fallback if Chunkr fails.

    Typical downstream tasks include:

    - Letting an LLM **summarise** or **answer questions about** a
    presentation, spreadsheet, contract, or research paper.  
    - Quickly **indexing** large document batches for semantic search.  
    - **Captioning media assets** (images & video) to improve accessibility.  
    - Turning "opaque" binary files into human‑readable text for diffing or
    version control.

    ### Args
    `workspace_dir` *(str)* – The unified workspace directory path where all files are located
    `filename` *(str)* – The filename within the workspace directory to process

    ### Returns
    `str` – On success, a **plain‑text or Markdown** representation of the
    file's meaningful content.  
    If the file type is unsupported or an extraction error occurs, the tool
    raises an exception containing a diagnostic message.
    """

    def __init__(self, cache_dir: Optional[str] = None):
        self.excel_tool = ExcelToolkit()
        self.cache_dir = cache_dir or "tmp/"

    # --------------------------------------------------------------------- #
    #  Public façade
    # --------------------------------------------------------------------- #
    @retry(tries=3, delay=1, backoff=2, exceptions=(Exception,))
    def extract_document_content(self, workspace_dir: str, filename: str) -> Tuple[bool, str]:
        """
        Extract content from a document located in the unified workspace.
        
        Args:
            workspace_dir: Path to the unified workspace directory
            filename: Name of the file to process (within workspace)
            
        Returns:
            Tuple[bool, str]: Success status and extracted content
        """
        # Build full document path using unified workspace structure
        document_path = _build_file_path(workspace_dir, filename)
        
        logger.debug(f"[extract_document_content] workspace_dir={workspace_dir}, filename={filename}, full_path={document_path}")

        # Verify file exists in workspace
        if not os.path.exists(document_path):
            logger.warning(f"File '{filename}' not found at path '{document_path}'.")
            return False, f"File '{filename}' not found."

        # 1. Images ----------------------------------------------------------------
        if document_path.lower().endswith((".jpg", ".jpeg", ".png")):
            caption = asyncio.run(
                ask_question_about_image(
                    workspace_dir,
                    filename,
                    "Please make a detailed caption about the image."
                )
            )
            return True, caption

        # 2. Audio -----------------------------------------------------------------
        if document_path.lower().endswith((".mp3", ".wav", ".m4a")):
            aai.settings.api_key = os.getenv("ASSEMBLYAI_API_KEY")
            config = aai.TranscriptionConfig(speech_model=aai.SpeechModel.best)
            transcript = aai.Transcriber(config=config).transcribe(document_path)
            logger.info(transcript.text)
            if transcript.status == "error":
                raise RuntimeError(f"Transcription failed: {transcript.error}")
            return True, transcript.text

        # 3. PPTX ------------------------------------------------------------------
        if document_path.lower().endswith(".pptx"):
            return True, asyncio.run(self._extract_pptx(document_path, workspace_dir))

        # 4. Spreadsheets -----------------------------------------------------------
        if document_path.lower().endswith((".xls", ".xlsx", ".csv")):
            return True, self.excel_tool.extract_excel_content(workspace_dir, filename)

        # 5. Zip --------------------------------------------------------------------
        if document_path.lower().endswith(".zip"):
            return True, f"The extracted files are: {self._unzip_file(document_path, workspace_dir)}"

        # 6. Simple text‑like formats ----------------------------------------------
        simple_readers = {
            ".py":  lambda p: open(p, encoding="utf‑8").read(),
            ".txt": lambda p: open(p, encoding="utf‑8").read(),
        }
        if any(document_path.lower().endswith(ext) for ext in simple_readers):
            reader = simple_readers[os.path.splitext(document_path)[1]]
            return True, reader(document_path)

        # 7. JSON                                                                   #
        if document_path.lower().endswith((".json", ".jsonl", ".jsonld")):
            return True, self._extract_json(document_path, encoding="utf‑8")
        

        # 8. XML                                                                    #
        if document_path.lower().endswith(".xml"):
            data = open(document_path, encoding="utf‑8").read()
            try:
                return True, xmltodict.parse(data)
            except Exception:
                return True, data

        # 9. DOCX → markdown -------------------------------------------------------
        if document_path.lower().endswith(".docx"):
            md_filename = f"{os.path.basename(filename)}.md"
            md_path = _build_file_path(workspace_dir, md_filename)
            docx_to_markdown(document_path, md_path)
            return True, open(md_path, encoding="utf‑8").read()

        # 10. MOV video ------------------------------------------------------------
        if document_path.lower().endswith(".mov"):
            description = ask_question_about_video(
                workspace_dir, filename, "Please make a detailed description about the video."
            )
            return True, description

        # 11. Fallback – Chunkr / PDF text -----------------------------------------
        return self._try_chunkr_then_fallback(document_path)

    # ------------------------------------------------------------------------- #
    #  helpers
    # ------------------------------------------------------------------------- #
    def _extract_json(self, json_path: str, encoding: str = "utf‑8") -> str:
        with open(json_path, 'r', encoding=encoding) as f:
            if json_path.lower().endswith((".json",".jsonld")):
                return json.load(f)  
            elif json_path.lower().endswith(".jsonl"):
                return [json.loads(line) for line in f]                    

    async def _extract_pptx(self, pptx_path: str, workspace_dir: str) -> str:
        """
        Extract content from PowerPoint file, saving images to workspace directory.
        
        Args:
            pptx_path: Full path to the PPTX file
            workspace_dir: Workspace directory for saving extracted images
            
        Returns:
            str: Extracted text and image descriptions
        """
        prs = Presentation(pptx_path)
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        out = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            txt = [f"Page {slide_idx}"]
            captions = []
            img_count = 0

            for shape_idx, shape in enumerate(slide.shapes):
                if getattr(shape, "text", "").strip():
                    txt.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    img = Image.open(io.BytesIO(shape.image.blob))
                    img_filename = f"{base}_slide_{slide_idx}_img_{shape_idx}.png"
                    img_path = _build_file_path(workspace_dir, img_filename)
                    img.save(img_path)
                    captions.append(
                        f"Image {img_count}: "
                        + await ask_question_about_image(
                            workspace_dir,
                            filename,
                            "Please make a detailed caption about the image."
                        )
                    )

            out.append("\n".join(txt + captions))

        return "\n\n".join(out)

    def _try_chunkr_then_fallback(self, path: str) -> Tuple[bool, str]:
        try:
            text = asyncio.run(
                self._extract_with_chunkr(path, output_format="markdown")
            )
            return True, text
        except Exception as e:
            logger.warning(f"Chunkr failed: {e}")
            if path.lower().endswith(".pdf"):
                try:
                    from PyPDF2 import PdfReader
                    text = "".join(
                        p.extract_text() for p in PdfReader(open(path, "rb")).pages
                    )
                    return True, text
                except Exception as e2:
                    return False, f"PDF fallback failed: {e2}"
            return False, f"Unsupported file type or processing error: {e}"

    async def _extract_with_chunkr(
        self, path: str, output_format: Literal["json", "markdown"] = "markdown"
    ) -> str:
        chunkr = Chunkr(api_key=os.getenv("CHUNKR_API_KEY"))
        result = await chunkr.upload(path)

        if result.status == "Failed":
            raise RuntimeError(result.message)

        out_path = f"{os.path.basename(path)}.{ 'json' if output_format=='json' else 'md' }"
        (result.json if output_format == "json" else result.markdown)(out_path)
        return open(out_path, encoding="utf‑8").read()

    def _unzip_file(self, zip_path: str, workspace_dir: str) -> List[str]:
        """
        Extract ZIP file contents to workspace directory.
        
        Args:
            zip_path: Path to the ZIP file
            workspace_dir: Target workspace directory for extraction
            
        Returns:
            List[str]: List of extracted filenames
        """
        import zipfile
        
        extracted_files = []
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for member in zip_ref.namelist():
                # Extract to workspace directory with safe filename
                safe_filename = os.path.basename(member)
                if safe_filename:  # Skip directory entries
                    extract_path = _build_file_path(workspace_dir, safe_filename)
                    with zip_ref.open(member) as source, open(extract_path, 'wb') as target:
                        target.write(source.read())
                    extracted_files.append(safe_filename)
        
        return extracted_files


# --------------------------------------------------------------------------- #
#  FastMCP server
# --------------------------------------------------------------------------- #
mcp = FastMCP("document_processing")
toolkit = DocumentProcessingToolkit()


@mcp.tool()
async def process_document(
    filename: str, 
    task_cache_dir: str | None = None
) -> str:
    """
    Process a document from the unified workspace directory and extract its content.
    
    **File Location Requirements**:
    - All files must be placed in the workspace directory (task_cache_dir/workspace/)
    - This ensures consistency with code execution and other tools
    - The agent_core.py manages the task_cache_dir and passes it to tools
    
    Args:
        filename (str): Name of the file to process (within workspace directory)
        task_cache_dir (str, optional): Task-specific cache directory path.
                                       If not provided, uses AGENT_CACHE_DIR environment variable.
    
    Returns:
        str: Extracted content as plain text or Markdown
        
    Example Usage:
        await process_document("report.pdf", "/path/to/task_cache")
        await process_document("presentation.pptx")
    """
    # Get unified workspace directory
    workspace_dir = _get_workspace_dir(task_cache_dir)
    
    # Initialize toolkit and process document
    toolkit = DocumentProcessingToolkit()
    
    try:
        success, content = await anyio.to_thread.run_sync(
            toolkit.extract_document_content, workspace_dir, filename
        )
        
        if success:
            return content
        else:
            raise ValueError(f"Failed to process document: {content}")
            
    except Exception as e:
        logger.error(f"Document processing failed for {filename}: {e}")
        raise ValueError(f"Document processing error: {str(e)}")


# --------------------------------------------------------------------------- #
#  Entrypoint
# --------------------------------------------------------------------------- #
if __name__ == "__main__":
    mcp.run(transport="stdio")
